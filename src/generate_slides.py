"""
Generate MultiCue-DF PowerPoint presentation (15 slides, 16:9, dark theme).
Usage: python src/generate_slides.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_ROOT = Path(__file__).resolve().parent.parent
PLOTS = PROJECT_ROOT / "plots"

# ── Slide geometry ────────────────────────────────────────────────────────────
SW  = Inches(13.333)   # 16:9 width
SH  = Inches(7.5)      # 16:9 height
M   = Inches(0.35)     # left/right margin
CT  = Inches(1.08)     # content top (below header)
CW  = SW - 2 * M      # usable content width  ≈ 12.63"

# ── Colour palette (red-black theme) ─────────────────────────────────────────
BG     = RGBColor(0x1a, 0x00, 0x00)   # very dark red-black background
HDR    = RGBColor(0x3a, 0x00, 0x00)   # dark red header / box fill
PANEL  = RGBColor(0x22, 0x00, 0x00)   # dark red panel
TEAL   = RGBColor(0xcc, 0x00, 0x00)   # primary red accent (replaces teal)
ORANGE = RGBColor(0xff, 0x44, 0x44)   # lighter red (replaces orange)
GREEN  = RGBColor(0xe6, 0x33, 0x33)   # medium red (replaces green)
DKGRN  = RGBColor(0x5a, 0x00, 0x00)   # very dark red (best-row bg)
LTGRN  = RGBColor(0xff, 0xcc, 0xcc)   # light pink-red (best-row text)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LGREY  = RGBColor(0xff, 0xcc, 0xcc)   # light pink-red (captions / subtext)
DGREY  = RGBColor(0xff, 0x99, 0x99)   # dim pink (faded text)
GOLD   = RGBColor(0xff, 0x66, 0x66)   # light red (replaces gold)

FONT = "Calibri"

# ── Ablation data (from results/ablation_summary.csv) ────────────────────────
ABLATION = [
    ("1",  "E01 — MLP Baseline",               "66.81", "0.718", "0.776", "77.2M"),
    ("2",  "E02 — Single CNN Baseline",         "83.24", "0.658", "0.906", "0.42M"),
    ("3",  "E03 — Stream 1 Only (ResNet-18)",   "95.10", "0.986", "0.971", "10.9M"),
    ("4",  "E04 — Stream 2 Only (Eye/Mouth)",   "85.71", "0.650", "0.923", "1.31M"),
    ("5",  "E05 — Stream 3 Only (FFT-CNN)",     "71.86", "0.552", "0.830", "0.14M"),
    ("6",  "E06 — Stream 1 + Stream 2",         "96.14", "0.990", "0.978", "12.2M"),
    ("7",  "E07 — Stream 1 + Stream 3",         "95.90", "0.985", "0.976", "11.0M"),
    ("8",  "E08 — Stream 2 + Stream 3",         "78.57", "0.623", "0.876", "1.45M"),
    ("9",  "E09 — Full Model (Adam, d=0.5)",    "95.48", "0.984", "0.974", "12.4M"),
    ("10", "E10 — Full Model (SGD)",            "92.76", "0.957", "0.958", "12.4M"),
    ("11", "E11 — Full Model (d=0.3)  ★ BEST", "96.57", "0.990", "0.980", "12.4M"),
    ("12", "E12 — Full Model (No Attn)",        "96.43", "0.992", "0.979", "12.3M"),
]
BEST_IDX = 10  # 0-based index of E11

# =============================================================================
# Low-level helpers
# =============================================================================

def new_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_rect(slide, l, t, w, h, fill_c=None, line_c=None, line_w=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_c:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_c
    else:
        shape.fill.background()
    shape.line.width = Pt(line_w)
    if line_w > 0 and line_c:
        shape.line.color.rgb = line_c
    elif line_w == 0:
        shape.line.fill.background()
    return shape


def set_run(run, text, size, color, bold=False, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT


def set_para(p, text, size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, space_before=0):
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    set_run(p.add_run(), text, size, color, bold, italic)


def add_line(tf, text, size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, sp=0):
    p = tf.add_paragraph()
    set_para(p, text, size, color, bold, italic, align, sp)
    return p


def txt1(slide, text, l, t, w, h, size=13, color=WHITE, bold=False,
         align=PP_ALIGN.LEFT, italic=False):
    """Single-paragraph text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf.paragraphs[0], text, size, color, bold, italic, align)
    return tb


def txtbox(slide, l, t, w, h):
    """Return (textbox, text_frame) with word wrap on."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def embed_img(slide, path, l, t, target_w):
    p = Path(path)
    if not p.exists():
        txt1(slide, f"[Missing: {p.name}]", l, t, target_w, Inches(0.4),
             size=9, color=ORANGE)
        return None, 0
    pic = slide.shapes.add_picture(str(p), l, t)
    scale = target_w / pic.width
    pic.width = int(target_w)
    pic.height = int(pic.height * scale)
    return pic, pic.height


def fmt_cell(cell, text, size=9, fg=WHITE, bg=None, bold=False,
             align=PP_ALIGN.CENTER):
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    if p.runs:
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = fg
        r.font.name = FONT


def header_bar(slide, title, num, total=15):
    """Consistent dark-blue header bar with teal left strip."""
    set_bg(slide)
    add_rect(slide, 0, 0, SW, Inches(1.0), fill_c=HDR)
    add_rect(slide, 0, 0, Inches(0.07), Inches(1.0), fill_c=TEAL)
    txt1(slide, title,
         Inches(0.22), Inches(0.12), SW - Inches(2.0), Inches(0.82),
         size=25, color=WHITE, bold=True)
    txt1(slide, f"{num} / {total}",
         SW - Inches(1.9), Inches(0.32), Inches(1.7), Inches(0.38),
         size=11, color=LGREY, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, SH - Inches(0.055), SW, Inches(0.055), fill_c=TEAL)


def bullet_box(slide, l, t, w, h, accent, label, body, body_size=10):
    """Colored panel with left accent strip, bold label, body text."""
    add_rect(slide, l, t, w, h, fill_c=PANEL)
    add_rect(slide, l, t, Inches(0.07), h, fill_c=accent)
    txt1(slide, label,
         l + Inches(0.17), t + Inches(0.09), w - Inches(0.22), Inches(0.36),
         size=11, color=accent, bold=True)
    _, tf = txtbox(slide, l + Inches(0.17), t + Inches(0.44),
                   w - Inches(0.22), h - Inches(0.52))
    set_para(tf.paragraphs[0], body, body_size, WHITE)


# =============================================================================
# Slides
# =============================================================================

def s01_title(prs):
    slide = new_blank(prs)
    set_bg(slide)

    # Top teal bar
    add_rect(slide, 0, 0, SW, Inches(0.07), fill_c=TEAL)

    # Title panel
    add_rect(slide, M, Inches(0.95), CW, Inches(2.8), fill_c=HDR)
    add_rect(slide, M, Inches(0.95), Inches(0.1), Inches(2.8), fill_c=TEAL)

    _, tf = txtbox(slide, M + Inches(0.25), Inches(1.1),
                   CW - Inches(0.4), Inches(2.4))
    set_para(tf.paragraphs[0], "MultiCue-DF", 42, WHITE, bold=True)
    add_line(tf, "Multi-Cue Deepfake Detection via Parallel Stream CNNs",
             18, TEAL, bold=False, sp=4)
    add_line(tf, "with Squeeze-and-Excitation Fusion  ·  FaceForensics++ C23",
             13, LGREY, italic=True, sp=2)

    # Four stat boxes
    stats = [("12.4M", "Parameters"), ("96.57%", "Test Accuracy"),
             ("0.990",  "AUC"),        ("3",      "Parallel Streams")]
    bw = (CW - Inches(0.3)) / 4
    for i, (val, lbl) in enumerate(stats):
        bx = M + i * (bw + Inches(0.1))
        add_rect(slide, bx, Inches(4.05), bw, Inches(1.25), fill_c=PANEL)
        add_rect(slide, bx, Inches(4.05), bw, Inches(0.065), fill_c=ORANGE)
        txt1(slide, val,
             bx, Inches(4.12), bw, Inches(0.7),
             size=30, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
        txt1(slide, lbl,
             bx, Inches(4.8), bw, Inches(0.4),
             size=10, color=LGREY, align=PP_ALIGN.CENTER)

    # Author block
    add_rect(slide, M, Inches(5.55), CW, Inches(1.05), fill_c=PANEL)
    add_rect(slide, M, Inches(5.55), Inches(0.07), Inches(1.05), fill_c=GOLD)
    _, tf2 = txtbox(slide, M + Inches(0.22), Inches(5.63), CW - Inches(0.35), Inches(0.9))
    set_para(tf2.paragraphs[0], "Humna Tariq", 17, WHITE, bold=True)
    add_line(tf2, "CS-419 Deep Learning  ·  SEECS, NUST  ·  May 2026",
             12, LGREY, sp=3)

    add_rect(slide, 0, SH - Inches(0.065), SW, Inches(0.065), fill_c=TEAL)


def s02_problem(prs):
    slide = new_blank(prs)
    header_bar(slide, "Why Deepfake Detection Matters", 2)

    # Left column — three threat panels
    lw = CW * 0.60
    panels = [
        (TEAL,   "SCALE",            "Modern GAN-based fakes are near-photorealistic and "
                                     "can be produced in minutes. Autoencoder swaps, neural "
                                     "texture rendering, and expression transfer are all "
                                     "open-source and widely accessible."),
        (ORANGE, "GLOBAL IMPACT",    "Deepfakes enable political manipulation, financial fraud, "
                                     "non-consensual imagery, and targeted harassment. A single "
                                     "viral clip can cause irreversible reputational damage."),
        (GREEN,  "BRITTLE DETECTORS","Single-stream methods fail: spatial detectors are fooled "
                                     "by expression-transfer fakes; frequency detectors are "
                                     "defeated by post-processing. No single cue is robust."),
    ]
    by = CT + Inches(0.18)
    for color, label, body in panels:
        bullet_box(slide, M, by, lw, Inches(1.6), color, label, body)
        by += Inches(1.75)

    # Right — Pakistan context box
    rx = M + lw + Inches(0.22)
    rw = CW - lw - Inches(0.22)
    add_rect(slide, rx, CT + Inches(0.18), rw, Inches(5.1), fill_c=HDR)
    add_rect(slide, rx, CT + Inches(0.18), rw, Inches(0.065), fill_c=ORANGE)
    txt1(slide, "PAKISTAN CONTEXT",
         rx + Inches(0.12), CT + Inches(0.3), rw - Inches(0.2), Inches(0.38),
         size=12, color=ORANGE, bold=True)
    pak = [
        "2024 elections: AI-generated videos of politicians "
        "circulated on WhatsApp before any fact-check was possible.",
        "Social velocity: content goes viral in minutes — far "
        "faster than any manual verification process.",
        "Infrastructure: low-bandwidth Android phones, "
        "WhatsApp re-encoding (C40-equivalent), Urdu/Pashto metadata.",
        "Targeted harassment: deepfake imagery weaponised against "
        "women journalists and political activists.",
    ]
    _, tf = txtbox(slide, rx + Inches(0.12), CT + Inches(0.75),
                   rw - Inches(0.2), Inches(4.1))
    set_para(tf.paragraphs[0], pak[0], 10, WHITE)
    for ln in pak[1:]:
        add_line(tf, "", 5, LGREY, sp=2)
        add_line(tf, ln, 10, WHITE, sp=0)

    # Bottom strip
    add_rect(slide, 0, SH - Inches(0.65), SW, Inches(0.58), fill_c=TEAL)
    txt1(slide,
         "Solution: fuse spatial semantics + periocular texture + frequency fingerprints "
         "— so no single failure mode can defeat the system",
         Inches(0.15), SH - Inches(0.6), SW - Inches(0.3), Inches(0.5),
         size=12, color=HDR, bold=True, align=PP_ALIGN.CENTER)


def s03_solution(prs):
    slide = new_blank(prs)
    header_bar(slide, "MultiCue-DF — Architecture Overview", 3)

    # Three stream boxes side by side
    streams = [
        ("STREAM 1", "Full Face  224×224 RGB",
         "ResNet-18\n(ImageNet pretrained)", "512-d", TEAL),
        ("STREAM 2", "Eye/Mouth  224×224 RGB",
         "Custom 4-Block\nCNN (from scratch)", "256-d", ORANGE),
        ("STREAM 3", "FFT Map  224×224",
         "Shallow 3-Block\nCNN + LeakyReLU", "128-d", GREEN),
    ]
    bw = Inches(3.4);  bh = Inches(3.7);  by = CT + Inches(0.25)
    for i, (label, inp, arch, out_dim, color) in enumerate(streams):
        bx = M + i * (bw + Inches(0.22))
        add_rect(slide, bx, by, bw, bh, fill_c=HDR)
        add_rect(slide, bx, by, bw, Inches(0.065), fill_c=color)
        txt1(slide, label,
             bx, by + Inches(0.13), bw, Inches(0.4),
             size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        txt1(slide, "Input:",
             bx + Inches(0.15), by + Inches(0.6), bw - Inches(0.3), Inches(0.28),
             size=9, color=LGREY)
        txt1(slide, inp,
             bx, by + Inches(0.88), bw, Inches(0.55),
             size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt1(slide, arch,
             bx, by + Inches(1.55), bw, Inches(0.7),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
        # output badge
        add_rect(slide, bx + Inches(0.55), by + bh - Inches(0.65),
                 bw - Inches(1.1), Inches(0.5), fill_c=color)
        txt1(slide, f"Output: {out_dim}",
             bx + Inches(0.55), by + bh - Inches(0.63),
             bw - Inches(1.1), Inches(0.46),
             size=12, color=HDR, bold=True, align=PP_ALIGN.CENTER)

    # Fusion box (right side)
    fx = M + 3 * (bw + Inches(0.22)) - Inches(0.05)
    fw = CW - (3 * (bw + Inches(0.22)) - Inches(0.05))
    add_rect(slide, fx, by, fw, bh, fill_c=PANEL)
    add_rect(slide, fx, by, fw, Inches(0.065), fill_c=GOLD)
    txt1(slide, "FUSION HEAD",
         fx, by + Inches(0.13), fw, Inches(0.4),
         size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    _, tf = txtbox(slide, fx + Inches(0.1), by + Inches(0.6),
                   fw - Inches(0.2), Inches(2.7))
    steps = ["Concat →", "896-d vector", "SE Attention", "f ⊙ s",
             "Linear(896→256)", "BN + ReLU + Dropout", "Linear(256→2)", "Real / Fake"]
    cols  = [LGREY, WHITE, TEAL, TEAL, WHITE, WHITE, WHITE, GREEN]
    bolds = [False, True, True, True, False, False, False, True]
    set_para(tf.paragraphs[0], steps[0], 10, cols[0], bolds[0], align=PP_ALIGN.CENTER)
    for step, c, b in zip(steps[1:], cols[1:], bolds[1:]):
        add_line(tf, step, 10, c, bold=b, align=PP_ALIGN.CENTER, sp=2)

    # Stats strip
    add_rect(slide, M, SH - Inches(1.12), CW, Inches(0.82), fill_c=PANEL)
    add_rect(slide, M, SH - Inches(1.12), CW, Inches(0.065), fill_c=GOLD)
    txt1(slide,
         "12,438,234 total parameters  ·  Best: 96.57% accuracy  ·  AUC 0.990  ·  "
         "F1 0.980  ·  FaceForensics++ C23",
         M + Inches(0.2), SH - Inches(1.02), CW - Inches(0.4), Inches(0.65),
         size=12, color=GOLD, bold=True, align=PP_ALIGN.CENTER)


def s04_dataset(prs):
    slide = new_blank(prs)
    header_bar(slide, "Dataset: FaceForensics++ C23", 4)

    # Manipulation method table
    tw = CW * 0.60
    methods = [
        ("Deepfakes",          "Identity swap via autoencoder trained per-person"),
        ("Face2Face",          "Expression transfer from source to target face"),
        ("FaceShifter",        "High-fidelity identity swap (occlusion-aware)"),
        ("FaceSwap",           "3D model-based face replacement"),
        ("NeuralTextures",     "Neural rendering of face texture patch"),
        ("DeepFakeDetection",  "Extended GAN set (Google / FaceForensics++ team)"),
    ]
    n = len(methods) + 1
    tshp = slide.shapes.add_table(n, 2, M, CT + Inches(0.15), tw, Inches(4.85))
    tbl = tshp.table
    tbl.columns[0].width = int(tw * 0.315)
    tbl.columns[1].width = int(tw * 0.685)
    rh = int(Inches(4.85) / n)
    for row in tbl.rows:
        row.height = rh
    fmt_cell(tbl.cell(0, 0), "Manipulation Method", 10, WHITE, HDR, bold=True)
    fmt_cell(tbl.cell(0, 1), "Description", 10, WHITE, HDR, bold=True, align=PP_ALIGN.LEFT)
    for i, (name, desc) in enumerate(methods, 1):
        bg_c = PANEL if i % 2 == 0 else None
        fmt_cell(tbl.cell(i, 0), name, 10, TEAL, bg_c, bold=True, align=PP_ALIGN.LEFT)
        fmt_cell(tbl.cell(i, 1), desc, 10, WHITE, bg_c, align=PP_ALIGN.LEFT)

    # Right: stat boxes
    rx = M + tw + Inches(0.25)
    rw = CW - tw - Inches(0.25)
    stats = [
        ("7,000",       "Total Videos",       TEAL),
        ("~14,000",     "Sampled Frames",      ORANGE),
        ("70 / 15 / 15","Train / Val / Test",  GREEN),
        ("224 × 224",   "Image Resolution",    GOLD),
        ("C23 H.264",   "Compression Level",   LGREY),
        ("Seed = 42",   "Reproducibility",     LGREY),
    ]
    bh = Inches(0.68)
    by = CT + Inches(0.15)
    for val, lbl, color in stats:
        add_rect(slide, rx, by, rw, bh, fill_c=PANEL)
        add_rect(slide, rx, by, Inches(0.06), bh, fill_c=color)
        txt1(slide, val,
             rx, by + Inches(0.04), rw, Inches(0.38),
             size=19, color=color, bold=True, align=PP_ALIGN.CENTER)
        txt1(slide, lbl,
             rx, by + Inches(0.42), rw, Inches(0.22),
             size=9, color=LGREY, align=PP_ALIGN.CENTER)
        by += bh + Inches(0.1)


def s05_architecture(prs):
    slide = new_blank(prs)
    header_bar(slide, "Three-Stream Architecture (Detailed)", 5)

    streams = [
        ("Stream 1 — Full Face",
         "ResNet-18 (ImageNet pretrained)",
         ["Input: 224×224 RGB full-face crop",
          "Layers 0–5: FROZEN (conv1 → layer2)",
          "  preserves universal low-level features",
          "Layers 3–4 + FC: fine-tuned",
          "FC: Linear(512→512) + ReLU + Dropout(0.3)",
          "Output: 512-d feature vector"],
         TEAL, "~11.2M params"),
        ("Stream 2 — Eye / Mouth Region",
         "Custom 4-Block CNN (trained from scratch)",
         ["Input: 224×224 RGB periocular crop",
          "B1: Conv→BN→ReLU×2 → 32ch + MaxPool",
          "B2: 32→64ch + MaxPool",
          "B3: 64→128ch + MaxPool",
          "B4: 128→256ch + AdaptiveAvgPool",
          "Head: Linear(256→256) + ReLU + Dropout(0.3)",
          "Output: 256-d feature vector"],
         ORANGE, "~1.31M params"),
        ("Stream 3 — FFT Frequency Domain",
         "Shallow 3-Block CNN (LeakyReLU α=0.1)",
         ["Input: 224×224 log FFT magnitude map",
          "  (zero-mean, unit-variance normalised)",
          "Conv→BN→LReLU×3: 3→32→64→128ch",
          "MaxPool × 2 + AdaptiveAvgPool",
          "Head: Linear(128→128) + LReLU + Dropout(0.3)",
          "Output: 128-d feature vector"],
         GREEN, "~145K params"),
    ]
    bw = (CW - Inches(0.3)) / 3
    bh = Inches(4.55)
    by = CT + Inches(0.2)
    for i, (title, subtitle, bullets, color, params) in enumerate(streams):
        bx = M + i * (bw + Inches(0.15))
        add_rect(slide, bx, by, bw, bh, fill_c=HDR)
        add_rect(slide, bx, by, bw, Inches(0.065), fill_c=color)
        txt1(slide, title,
             bx + Inches(0.1), by + Inches(0.12), bw - Inches(0.2), Inches(0.38),
             size=11, color=color, bold=True)
        txt1(slide, subtitle,
             bx + Inches(0.1), by + Inches(0.52), bw - Inches(0.2), Inches(0.32),
             size=9, color=LGREY, italic=True)
        add_rect(slide, bx + Inches(0.1), by + Inches(0.88),
                 bw - Inches(0.2), Inches(0.02), fill_c=color)
        _, tf = txtbox(slide, bx + Inches(0.1), by + Inches(0.97),
                       bw - Inches(0.2), bh - Inches(1.45))
        set_para(tf.paragraphs[0], bullets[0], 9, WHITE)
        for b in bullets[1:]:
            c = LGREY if b.startswith("  ") else WHITE
            add_line(tf, b, 9, c, sp=3)
        add_rect(slide, bx + Inches(0.12), by + bh - Inches(0.5),
                 bw - Inches(0.24), Inches(0.42), fill_c=color)
        txt1(slide, params,
             bx + Inches(0.12), by + bh - Inches(0.48),
             bw - Inches(0.24), Inches(0.38),
             size=10, color=HDR, bold=True, align=PP_ALIGN.CENTER)

    # Fusion formula strip
    add_rect(slide, 0, SH - Inches(0.68), SW, Inches(0.6), fill_c=PANEL)
    add_rect(slide, 0, SH - Inches(0.68), SW, Inches(0.05), fill_c=GOLD)
    txt1(slide,
         "FUSION:  cat([512, 256, 128]) = 896-d  →  s = σ(W₂·ReLU(W₁·f))  "
         "→  f ⊙ s  →  Linear(896→256) → BN1d → ReLU → Dropout → Linear(256→2)",
         Inches(0.15), SH - Inches(0.62), SW - Inches(0.3), Inches(0.52),
         size=11, color=GOLD, bold=True, align=PP_ALIGN.CENTER)


def s06_ablation_table(prs):
    slide = new_blank(prs)
    header_bar(slide, "12 Controlled Ablation Experiments", 6)

    n_rows = len(ABLATION) + 1
    n_cols = 6
    th = SH - CT - Inches(0.35)
    tshp = slide.shapes.add_table(n_rows, n_cols, M, CT + Inches(0.08), CW, th)
    tbl = tshp.table

    cw_list = [int(CW * p) for p in [0.044, 0.448, 0.15, 0.126, 0.126, 0.106]]
    for j, w in enumerate(cw_list):
        tbl.columns[j].width = w
    rh = int(th / n_rows)
    for row in tbl.rows:
        row.height = rh

    for j, h in enumerate(["#", "Experiment", "Acc (%)", "AUC", "F1", "Params"]):
        fmt_cell(tbl.cell(0, j), h, 10, WHITE, HDR, bold=True)

    for i, (num, name, acc, auc, f1, params) in enumerate(ABLATION):
        ri = i + 1
        is_best = (i == BEST_IDX)
        bg_c = (DKGRN if is_best else (PANEL if i % 2 == 0 else None))
        fg_c = LTGRN if is_best else WHITE
        fmt_cell(tbl.cell(ri, 0), num,    8.5, fg_c, bg_c)
        fmt_cell(tbl.cell(ri, 1), name,   8.5, fg_c, bg_c, bold=is_best, align=PP_ALIGN.LEFT)
        fmt_cell(tbl.cell(ri, 2), acc,    8.5, fg_c, bg_c, bold=is_best)
        fmt_cell(tbl.cell(ri, 3), auc,    8.5, fg_c, bg_c)
        fmt_cell(tbl.cell(ri, 4), f1,     8.5, fg_c, bg_c)
        fmt_cell(tbl.cell(ri, 5), params, 8.5, fg_c, bg_c)


def s07_ablation_chart(prs):
    slide = new_blank(prs)
    header_bar(slide, "Ablation Study — Visual Comparison", 7)
    embed_img(slide, PLOTS / "ablation_comparison.png",
              M, CT + Inches(0.08), CW)
    txt1(slide,
         "All 12 experiments sorted by accuracy.  Best: E11 (dropout=0.3) — 96.57% accuracy",
         M, SH - Inches(0.38), CW, Inches(0.3),
         size=10, color=LGREY, italic=True, align=PP_ALIGN.CENTER)


def s08_ablation_findings(prs):
    slide = new_blank(prs)
    header_bar(slide, "What the Ablation Tells Us", 8)

    boxes = [
        (TEAL,   "1.  Transfer Learning Dominates",
                 "ResNet-18 alone (E03):  95.10%\n"
                 "Custom CNN from scratch (E04):  85.71%\n"
                 "Gap: +9.4pp — ImageNet features generalise strongly "
                 "to face manipulation detection at 14K sample scale."),
        (ORANGE, "2.  Adam >> SGD  (+2.7pp)",
                 "Adam (E09): 95.48%   vs   SGD (E10): 92.76%\n"
                 "Three parallel streams create heterogeneous gradient scales. "
                 "Adam's adaptive per-parameter LR handles this; "
                 "SGD's single LR cannot balance a frozen ResNet + two scratch CNNs."),
        (GREEN,  "3.  Dropout 0.3 > 0.5  (+1.09pp)",
                 "E09 (d=0.5): 95.48%   →   E11 (d=0.3): 96.57%\n"
                 "Model is in high-bias regime at d=0.5 — over-regularised. "
                 "Relaxing to 0.3 unlocks the full capacity of the 896-d fusion vector."),
        (GOLD,   "4.  FFT Stream Adds Noise via Simple Concat",
                 "Stream1+2 (E06): 96.14%  >  Full model (E09): 95.48%\n"
                 "Adding Stream 3 (FFT, standalone 71.9%) HURTS accuracy. "
                 "Weak frequency features introduce noise. Needs gated fusion, "
                 "not naive concatenation."),
    ]
    bw = (CW - Inches(0.2)) / 2
    bh = Inches(2.58)
    for i, (color, label, body) in enumerate(boxes):
        bx = M + (i % 2) * (bw + Inches(0.2))
        by = CT + Inches(0.18) + (i // 2) * (bh + Inches(0.18))
        add_rect(slide, bx, by, bw, bh, fill_c=HDR)
        add_rect(slide, bx, by, Inches(0.07), bh, fill_c=color)
        txt1(slide, label,
             bx + Inches(0.17), by + Inches(0.1), bw - Inches(0.22), Inches(0.4),
             size=12, color=color, bold=True)
        _, tf = txtbox(slide, bx + Inches(0.17), by + Inches(0.56),
                       bw - Inches(0.22), bh - Inches(0.68))
        lines = body.split("\n")
        set_para(tf.paragraphs[0], lines[0], 10, WHITE)
        for ln in lines[1:]:
            add_line(tf, ln, 10, WHITE, sp=3)


def s09_curves(prs):
    slide = new_blank(prs)
    header_bar(slide, "Training Convergence — Top 5 Experiments", 9)
    embed_img(slide, PLOTS / "training_curves.png",
              M, CT + Inches(0.1), CW)
    txt1(slide,
         "Top 5 by test accuracy: E11, E12, E06, E07, E09.  "
         "All converge by epoch 15–18.  "
         "CosineAnnealingLR prevents loss spikes.  E11 achieves best final val acc.",
         M, SH - Inches(0.42), CW, Inches(0.32),
         size=10, color=LGREY, italic=True, align=PP_ALIGN.CENTER)


def s10_activation(prs):
    slide = new_blank(prs)
    header_bar(slide, "Activation Function Analysis — Stream 2 CNN", 10)

    # Image in top portion
    pic, ph = embed_img(slide, PLOTS / "activation_comparison.png",
                        M, CT + Inches(0.08), CW)

    # Three finding rows below
    findings = [
        (TEAL,
         "ELU WINS across all metrics",
         "ELU: 95.67% acc, AUC 0.9888, F1 0.9745  —  vs  "
         "ReLU: 94.57% / 0.9773 / 0.9682  and  Leaky ReLU: 94.29% / 0.9790 / 0.9663"),
        (ORANGE,
         "ELU also fastest: −13% training time",
         "162 min (ELU)  vs  187 min (ReLU)  vs  175 min (Leaky ReLU).  "
         "ELU's self-normalising property (mean closer to zero) speeds convergence."),
        (GREEN,
         "Dying ReLU in periocular CNN — root cause",
         "Manipulation artefacts in eye/mouth region produce near-zero "
         "pre-activations. ReLU's hard-zero threshold kills these neurons. "
         "ELU's smooth exponential (e^x−1 for x<0) preserves gradient flow."),
    ]
    bh = Inches(0.75)
    by = CT + Inches(3.65)
    for color, label, body in findings:
        add_rect(slide, M, by, CW, bh, fill_c=PANEL)
        add_rect(slide, M, by, Inches(0.07), bh, fill_c=color)
        txt1(slide, label,
             M + Inches(0.17), by + Inches(0.07), CW * 0.32, Inches(0.35),
             size=10, color=color, bold=True)
        _, tf = txtbox(slide, M + CW * 0.33, by + Inches(0.1),
                       CW * 0.67 - Inches(0.1), bh - Inches(0.15))
        set_para(tf.paragraphs[0], body, 10, WHITE)
        by += bh + Inches(0.07)


def s11_batchnorm(prs):
    slide = new_blank(prs)
    header_bar(slide, "Batch Normalisation Analysis — Stream 2 CNN", 11)

    embed_img(slide, PLOTS / "batchnorm_comparison.png",
              M, CT + Inches(0.08), CW)

    findings = [
        (ORANGE,
         "WITHOUT BN WINS: +0.38pp acc, +0.32 AUC, saves 31% training time",
         "With BN: 94.57% / 175.3 min   vs   Without BN: 94.95% / 121.5 min  "
         "→  53.9 minutes saved per run. Removing BN improves ALL metrics simultaneously."),
        (TEAL,
         "WHY BN IS REDUNDANT in Stream 2",
         "Stream 2 is sandwiched: ResNet-18 (S1) applies BN throughout its blocks; "
         "the fusion head applies BN1d on the 896-d concatenated vector. "
         "Per-stream BN inside S2 adds zero benefit — the signal is normalised at both ends."),
        (GREEN,
         "DESIGN PRINCIPLE for multi-stream networks",
         "Place BN at the fusion point, not inside each stream independently. "
         "Fusion-level normalisation handles inter-stream scale differences more efficiently "
         "than per-stream BN. Aligns with internal covariate shift theory (Ioffe & Szegedy 2015)."),
    ]
    bh = Inches(0.73)
    by = CT + Inches(3.6)
    for color, label, body in findings:
        add_rect(slide, M, by, CW, bh, fill_c=PANEL)
        add_rect(slide, M, by, Inches(0.07), bh, fill_c=color)
        txt1(slide, label,
             M + Inches(0.17), by + Inches(0.07), CW * 0.35, Inches(0.38),
             size=10, color=color, bold=True)
        _, tf = txtbox(slide, M + CW * 0.36, by + Inches(0.1),
                       CW * 0.64 - Inches(0.1), bh - Inches(0.15))
        set_para(tf.paragraphs[0], body, 10, WHITE)
        by += bh + Inches(0.07)


def s12_best_model(prs):
    slide = new_blank(prs)
    header_bar(slide, "Best Model — E11: Full Model (dropout=0.3)", 12)

    # Three metric boxes
    metrics = [("96.57%", "Test Accuracy", TEAL),
               ("0.9904",  "Test AUC",     ORANGE),
               ("0.9799",  "Test F1",      GREEN)]
    mw = (CW - Inches(0.3)) / 3
    for i, (val, lbl, color) in enumerate(metrics):
        bx = M + i * (mw + Inches(0.15))
        add_rect(slide, bx, CT + Inches(0.08), mw, Inches(1.05), fill_c=PANEL)
        add_rect(slide, bx, CT + Inches(0.08), mw, Inches(0.06), fill_c=color)
        txt1(slide, val,
             bx, CT + Inches(0.15), mw, Inches(0.62),
             size=32, color=color, bold=True, align=PP_ALIGN.CENTER)
        txt1(slide, lbl,
             bx, CT + Inches(0.76), mw, Inches(0.3),
             size=11, color=LGREY, align=PP_ALIGN.CENTER)

    # Confusion matrix image
    img_w = CW * 0.73
    embed_img(slide, PLOTS / "confusion_matrix.png",
              M, CT + Inches(1.25), img_w)

    # Right: summary notes
    rx = M + img_w + Inches(0.22)
    rw = CW - img_w - Inches(0.22)
    notes = [("2,100", "Test Samples"), ("~72", "Misclassified"),
             ("96.86%", "Best Val Acc\n(epoch 19)"), ("12.4M", "Parameters")]
    ny = CT + Inches(1.3)
    for val, lbl in notes:
        add_rect(slide, rx, ny, rw, Inches(1.1), fill_c=HDR)
        txt1(slide, val,
             rx, ny + Inches(0.1), rw, Inches(0.58),
             size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        txt1(slide, lbl,
             rx, ny + Inches(0.66), rw, Inches(0.38),
             size=9, color=LGREY, align=PP_ALIGN.CENTER)
        ny += Inches(1.18)

    txt1(slide,
         "Checkpoint: epoch 19, val acc 96.86%.  "
         "Balanced errors: ~36 FP + ~36 FN.  "
         "High AUC (0.990) confirms strong class separation.",
         M, SH - Inches(0.38), CW, Inches(0.3),
         size=9, color=DGREY, italic=True, align=PP_ALIGN.CENTER)


def s13_findings(prs):
    slide = new_blank(prs)
    header_bar(slide, "Key Findings — Summary", 13)

    rows = [
        (TEAL,   "1", "Multi-stream fusion outperforms all single-stream baselines  "
                      "(96.57% vs best single 95.10% — Stream 1 ResNet-18)"),
        (ORANGE, "2", "Transfer learning dominates at dataset scale:  "
                      "ResNet-18 (95.10%)  vs  custom CNN scratch (85.71%)  →  +9.4pp"),
        (GREEN,  "3", "Adam + CosineAnnealingLR is essential:  "
                      "Adam 95.48%  vs  SGD 92.76%  (−2.7pp for heterogeneous multi-stream grads)"),
        (GOLD,   "4", "Dropout tuning (0.5 → 0.3) gives +1.09pp — larger than adding SE attention "
                      "(+0.14pp). Model was over-regularised."),
        (TEAL,   "5", "ELU activation: +1.09pp vs ReLU and 13% faster training. "
                      "Smooth negative half-plane prevents dying neurons in periocular CNN."),
        (ORANGE, "6", "No BN in Stream 2: +0.38pp and −31% training time (−54 min/run). "
                      "Fusion BN1d makes per-stream BN redundant."),
        (GREEN,  "7", "FFT stream adds noise via naive concat (S1+S2 96.14% > Full 95.48%). "
                      "Needs gated fusion — not simple concatenation."),
    ]
    bh = Inches(0.7)
    by = CT + Inches(0.18)
    for color, num, text in rows:
        add_rect(slide, M, by, CW, bh, fill_c=PANEL)
        add_rect(slide, M, by, Inches(0.07), bh, fill_c=color)
        txt1(slide, num,
             M + Inches(0.15), by + Inches(0.15), Inches(0.32), Inches(0.4),
             size=16, color=color, bold=True)
        _, tf = txtbox(slide, M + Inches(0.5), by + Inches(0.12),
                       CW - Inches(0.58), bh - Inches(0.18))
        set_para(tf.paragraphs[0], text, 11, WHITE)
        by += bh + Inches(0.08)


def s14_realworld(prs):
    slide = new_blank(prs)
    header_bar(slide, "Deployment in Pakistan's Context", 14)

    col_w = (CW - Inches(0.25)) / 3
    col_h = SH - CT - Inches(0.4)

    sections = [
        ("DEPLOYMENT CHALLENGES", ORANGE, [
            ("Low-resource devices",
             "2–4 GB RAM Android, no GPU. Need MobileNetV3 distillation (<5M params)."),
            ("WhatsApp C40 compression",
             "Re-encoding destroys FFT fingerprints. Requires C40 fine-tuning."),
            ("Urdu/Pashto metadata",
             "Low-resource NLP limits multimodal text+video detection."),
        ]),
        ("ETHICAL CONSIDERATIONS", TEAL, [
            ("False positive impact",
             "3.4% of real videos wrongly flagged. Never use binary verdict — "
             "show confidence score + human review for borderline cases."),
            ("Dual-use tension",
             "Publishing failure modes helps adversaries. Responsible disclosure only."),
            ("Governance first",
             "Legal frameworks and media literacy may matter more than accuracy gains."),
        ]),
        ("INTEGRATION PATHWAYS", GREEN, [
            ("Browser extension",
             "For Dawn, Geo, ARY News: overlay confidence badge on video content."),
            ("WhatsApp Business API",
             "Server-side detection: append metadata flag before delivery in groups."),
            ("PEMRA broadcast monitoring",
             "Real-time deepfake flagging in Pakistan broadcast regulatory infrastructure."),
        ]),
    ]
    for i, (title, color, items) in enumerate(sections):
        bx = M + i * (col_w + Inches(0.125))
        add_rect(slide, bx, CT + Inches(0.1), col_w, col_h, fill_c=HDR)
        add_rect(slide, bx, CT + Inches(0.1), col_w, Inches(0.065), fill_c=color)
        txt1(slide, title,
             bx + Inches(0.1), CT + Inches(0.2), col_w - Inches(0.15), Inches(0.38),
             size=10, color=color, bold=True)
        by = CT + Inches(0.65)
        for subhdr, body in items:
            add_rect(slide, bx + Inches(0.08), by,
                     col_w - Inches(0.16), Inches(0.02), fill_c=color)
            txt1(slide, subhdr,
                 bx + Inches(0.1), by + Inches(0.07), col_w - Inches(0.15), Inches(0.3),
                 size=10, color=color, bold=True)
            _, tf = txtbox(slide, bx + Inches(0.1), by + Inches(0.4),
                           col_w - Inches(0.15), Inches(1.05))
            set_para(tf.paragraphs[0], body, 9, WHITE)
            by += Inches(1.6)


def s15_conclusion(prs):
    slide = new_blank(prs)
    header_bar(slide, "Conclusion & Future Work", 15)

    lw = CW * 0.50

    # Left: Conclusions header
    add_rect(slide, M, CT + Inches(0.1), lw, Inches(0.38), fill_c=TEAL)
    txt1(slide, "CONCLUSIONS",
         M + Inches(0.1), CT + Inches(0.15), lw - Inches(0.15), Inches(0.28),
         size=12, color=HDR, bold=True)

    conclusions = [
        ("96.57% accuracy", "competitive with SOTA using 12.4M params on FaceForensics++ C23"),
        ("Multi-stream fusion", "confirms three complementary cues outperform any single stream"),
        ("ELU + No BN in S2", "two best architectural improvements from targeted analysis"),
        ("Pakistan deployment", "requires distillation, C40 fine-tuning, and human-in-loop review"),
    ]
    by = CT + Inches(0.55)
    for val, body in conclusions:
        add_rect(slide, M, by, lw, Inches(1.05), fill_c=PANEL)
        add_rect(slide, M, by, Inches(0.06), Inches(1.05), fill_c=TEAL)
        txt1(slide, val,
             M + Inches(0.15), by + Inches(0.08), lw - Inches(0.22), Inches(0.35),
             size=11, color=TEAL, bold=True)
        _, tf = txtbox(slide, M + Inches(0.15), by + Inches(0.42),
                       lw - Inches(0.22), Inches(0.55))
        set_para(tf.paragraphs[0], body, 10, WHITE)
        by += Inches(1.13)

    # Right: Future work header
    rx = M + lw + Inches(0.22)
    rw = CW - lw - Inches(0.22)
    add_rect(slide, rx, CT + Inches(0.1), rw, Inches(0.38), fill_c=ORANGE)
    txt1(slide, "FUTURE WORK",
         rx + Inches(0.1), CT + Inches(0.15), rw - Inches(0.15), Inches(0.28),
         size=12, color=HDR, bold=True)

    future = [
        (ORANGE, "Gated Fusion for Stream 3",
                 "Learned gate suppresses weak FFT contribution when confidence "
                 "is low — fixes the S1+S2 > Full model counter-intuitive result"),
        (TEAL,   "Knowledge Distillation",
                 "Compress 12.4M → MobileNetV3 (<5M params) for on-device "
                 "inference on Pakistani low-end smartphones"),
        (GREEN,  "Temporal Modelling",
                 "Transformer / LSTM over 10-frame sequences for motion "
                 "inconsistencies + audio-lip sync deepfake detection"),
        (GOLD,   "Cross-Dataset Generalisation",
                 "Celeb-DF v2, DFDC, custom South Asian dataset — "
                 "quantify OOD performance and regional bias"),
    ]
    by = CT + Inches(0.55)
    for color, title, body in future:
        add_rect(slide, rx, by, rw, Inches(1.3), fill_c=HDR)
        add_rect(slide, rx, by, Inches(0.06), Inches(1.3), fill_c=color)
        txt1(slide, title,
             rx + Inches(0.15), by + Inches(0.08), rw - Inches(0.22), Inches(0.35),
             size=11, color=color, bold=True)
        _, tf = txtbox(slide, rx + Inches(0.15), by + Inches(0.44),
                       rw - Inches(0.22), Inches(0.77))
        set_para(tf.paragraphs[0], body, 10, WHITE)
        by += Inches(1.38)

    # Thank you footer
    add_rect(slide, 0, SH - Inches(0.72), SW, Inches(0.65), fill_c=HDR)
    add_rect(slide, 0, SH - Inches(0.72), SW, Inches(0.05), fill_c=TEAL)
    txt1(slide, "Thank you — Questions welcome",
         M, SH - Inches(0.65), CW * 0.55, Inches(0.55),
         size=18, color=WHITE, bold=True)
    txt1(slide, "humnatariq03@gmail.com  ·  CS-419 Deep Learning, SEECS NUST",
         M + CW * 0.55, SH - Inches(0.58), CW * 0.45, Inches(0.4),
         size=10, color=LGREY, align=PP_ALIGN.RIGHT)


# =============================================================================
# Main
# =============================================================================

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    s01_title(prs)
    s02_problem(prs)
    s03_solution(prs)
    s04_dataset(prs)
    s05_architecture(prs)
    s06_ablation_table(prs)
    s07_ablation_chart(prs)
    s08_ablation_findings(prs)
    s09_curves(prs)
    s10_activation(prs)
    s11_batchnorm(prs)
    s12_best_model(prs)
    s13_findings(prs)
    s14_realworld(prs)
    s15_conclusion(prs)

    out = PROJECT_ROOT / "presentations" / "MultiCue_DF_Slides.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    size_kb = out.stat().st_size // 1024
    print(f"Saved -> {out}  ({size_kb} KB,  15 slides)")


if __name__ == "__main__":
    main()
